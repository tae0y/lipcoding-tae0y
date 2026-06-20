"""Reusable themed deck builder on top of python-pptx.

Keeps every slide on one visual theme so build scripts stay declarative.
Run build scripts with:  uv run --with python-pptx python <script>.py
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

# 16:9 canvas (EMU). 13.333in x 7.5in.
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

MARGIN = Emu(685800)  # 0.75in

# CJK-safe font that ships on macOS/Windows.
FONT = "Malgun Gothic"
FONT_FALLBACK = "Apple SD Gothic Neo"

INK = RGBColor(0x1F, 0x29, 0x37)      # near-black slate
MUTED = RGBColor(0x6B, 0x72, 0x80)    # gray
LIGHT = RGBColor(0xF8, 0xFA, 0xFC)    # near-white bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _hex(color: str) -> RGBColor:
    return RGBColor.from_string(color.lstrip("#").upper())


def _set_font(run, *, size, bold=False, color=INK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT
    # East-Asian font hint
    rPr = run._r.get_or_add_rPr()
    import copy

    from pptx.oxml.ns import qn

    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", FONT)


class Deck:
    def __init__(self, title: str = "Deck", accent: str = "#2563EB"):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.accent = _hex(accent)
        self._blank = self.prs.slide_layouts[6]  # fully blank layout

    # --- internals ---------------------------------------------------
    def _slide(self, bg: RGBColor = LIGHT):
        s = self.prs.slides.add_slide(self._blank)
        fill = s.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
        return s

    def _accent_bar(self, slide, y=Emu(548640), h=Emu(91440)):
        bar = slide.shapes.add_shape(
            1, MARGIN, y, Emu(1280160), h  # MSO_SHAPE.RECTANGLE = 1
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.accent
        bar.line.fill.background()
        bar.shadow.inherit = False
        return bar

    def _textbox(self, slide, x, y, w, h):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        return tb, tf

    def _heading(self, slide, text):
        self._accent_bar(slide)
        tb, tf = self._textbox(slide, MARGIN, Emu(731520), Emu(10820400), Emu(960120))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = text
        _set_font(r, size=34, bold=True, color=INK)

    # --- public slide types ------------------------------------------
    def title_slide(self, title: str, subtitle: str = "", footer: str = ""):
        s = self._slide(bg=INK)
        # accent block
        block = s.shapes.add_shape(1, Emu(0), Emu(0), Emu(228600), SLIDE_H)
        block.fill.solid()
        block.fill.fore_color.rgb = self.accent
        block.line.fill.background()
        block.shadow.inherit = False

        tb, tf = self._textbox(s, MARGIN, Emu(2057400), Emu(10820400), Emu(2400000))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        _set_font(r, size=54, bold=True, color=WHITE)
        if subtitle:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(14)
            r2 = p2.add_run()
            r2.text = subtitle
            _set_font(r2, size=24, bold=False, color=_hex("#CBD5E1"))
        if footer:
            tb2, tf2 = self._textbox(
                s, MARGIN, Emu(6126480), Emu(10820400), Emu(548640)
            )
            rp = tf2.paragraphs[0].add_run()
            rp.text = footer
            _set_font(rp, size=14, color=MUTED)
        return s

    def section_slide(self, number: str, title: str):
        s = self._slide(bg=INK)
        tb, tf = self._textbox(s, MARGIN, Emu(2300000), Emu(10820400), Emu(2000000))
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = number
        _set_font(r, size=28, bold=True, color=self.accent)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        r2 = p2.add_run()
        r2.text = title
        _set_font(r2, size=48, bold=True, color=WHITE)
        return s

    def bullets_slide(self, title: str, bullets: list):
        s = self._slide()
        self._heading(s, title)
        tb, tf = self._textbox(
            s, MARGIN, Emu(1965960), Emu(10820400), Emu(4343400)
        )
        for i, item in enumerate(bullets):
            text, level = (item if isinstance(item, tuple) else (item, 0))
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            p.level = level
            bullet = "•   " if level == 0 else "–   "
            rb = p.add_run()
            rb.text = bullet
            _set_font(rb, size=20, bold=True, color=self.accent)
            r = p.add_run()
            r.text = text
            _set_font(r, size=20, bold=(level == 0), color=INK if level == 0 else MUTED)
        return s

    def two_column_slide(
        self, title, left_title, left, right_title, right
    ):
        s = self._slide()
        self._heading(s, title)
        col_w = Emu(5212080)
        gap = Emu(396240)
        self._column(s, MARGIN, col_w, left_title, left)
        self._column(s, MARGIN + col_w + gap, col_w, right_title, right)
        return s

    def _column(self, slide, x, w, col_title, items):
        # header chip
        chip = slide.shapes.add_shape(1, x, Emu(1965960), w, Emu(548640))
        chip.fill.solid()
        chip.fill.fore_color.rgb = self.accent
        chip.line.fill.background()
        chip.shadow.inherit = False
        ctf = chip.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = col_title
        _set_font(cr, size=18, bold=True, color=WHITE)

        tb, tf = self._textbox(slide, x, Emu(2697480), w, Emu(3500000))
        for i, item in enumerate(items):
            text, level = (item if isinstance(item, tuple) else (item, 0))
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            p.level = level
            rb = p.add_run()
            rb.text = "•   " if level == 0 else "–   "
            _set_font(rb, size=16, bold=True, color=self.accent)
            r = p.add_run()
            r.text = text
            _set_font(r, size=16, color=INK if level == 0 else MUTED)

    def table_slide(self, title, headers, rows):
        s = self._slide()
        self._heading(s, title)
        nrows = len(rows) + 1
        ncols = len(headers)
        gtbl = s.shapes.add_table(
            nrows, ncols, MARGIN, Emu(1965960), Emu(10820400), Emu(600000 * nrows)
        ).table
        for c, h in enumerate(headers):
            cell = gtbl.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.accent
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = h
            _set_font(run, size=16, bold=True, color=WHITE)
        for r_i, row in enumerate(rows, start=1):
            for c_i, val in enumerate(row):
                cell = gtbl.cell(r_i, c_i)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r_i % 2 else LIGHT
                para = cell.text_frame.paragraphs[0]
                run = para.add_run()
                run.text = str(val)
                _set_font(run, size=14, color=INK)
        return s

    def big_statement_slide(self, statement, sub=""):
        s = self._slide(bg=LIGHT)
        self._accent_bar(s, y=Emu(2880360))
        tb, tf = self._textbox(s, MARGIN, Emu(2880360), Emu(10820400), Emu(1800000))
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.space_before = Pt(18)
        r = p.add_run()
        r.text = statement
        _set_font(r, size=40, bold=True, color=INK)
        if sub:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(12)
            r2 = p2.add_run()
            r2.text = sub
            _set_font(r2, size=22, color=MUTED)
        return s

    def save(self, path: str):
        import os

        os.makedirs(os.path.dirname(path), exist_ok=True)
        self.prs.save(path)
        return path
